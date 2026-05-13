"""
Garment Catalog PowerPoint Generator

Layouts mirror the web catalog grouped style card:
  - Header: style id + garment type
  - Compact spec strip (REF / COMP / GSM / DATE) when data exists
  - Two-column image grid with object-cover–style crops (center crop via PIL)
  - Cover: hero image from the first style when available
  - Closing: brand footer + contact line
"""

from __future__ import annotations

import io
from pathlib import Path

from PIL import Image as PILImage
from PIL import ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Dimensions (16:9 widescreen) ──────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.500)

EMU_PER_INCH = 914400

# Colour tokens — light catalog (print-friendly), hierarchy similar to UI
PAGE_BG = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x18, 0x27)  # slate-900
SUBTITLE = RGBColor(0x37, 0x41, 0x51)  # slate-600
MUTED = RGBColor(0x6B, 0x72, 0x80)  # slate-500
FAINT = RGBColor(0x9C, 0xA3, 0xAF)  # slate-400
RULE = RGBColor(0xE5, 0xE7, 0xEB)  # slate-200
PANEL_BG = RGBColor(0xF3, 0xF4, 0xF6)  # gray-100
PANEL_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
COVER_PANEL = RGBColor(0xEE, 0xF0, 0xF4)
ACCENT = RGBColor(0x4F, 0x46, 0xE5)  # indigo-ish, echoes UI accent


def _emu(e) -> int:
    return int(e)


def _add_picture_cover(
    slide,
    img_path: str,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    dpi: int = 144,
) -> None:
    """Place image scaled/cropped to fill the slot (CSS object-cover)."""
    try:
        im = PILImage.open(img_path).convert("RGB")
        w_px = max(2, _emu(width) * dpi // EMU_PER_INCH)
        h_px = max(2, _emu(height) * dpi // EMU_PER_INCH)
        fitted = ImageOps.fit(im, (w_px, h_px), method=PILImage.Resampling.LANCZOS)
        buf = io.BytesIO()
        fitted.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, left, top, width=width, height=height)
    except Exception:
        _place_placeholder(slide, left, top, width, height, "IMG")


def generate_catalog_ppt(
    groups: list[dict],
    output_path: str,
    brand_name: str = "GARMENT COLLECTION",
) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_path = _first_cover_image(groups)
    _cover(prs, brand_name, cover_path)

    for i, grp in enumerate(groups, start=1):
        _product_slide(prs, grp, i)

    _closing(prs, brand_name)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _first_cover_image(groups: list[dict]) -> str | None:
    if not groups:
        return None
    first = groups[0]
    for img in first.get("images") or []:
        if img.get("image_type") != "front":
            continue
        p = img.get("processed_path") or img.get("original_path") or ""
        if p and Path(p).exists():
            return p
    for img in first.get("images") or []:
        p = img.get("processed_path") or img.get("original_path") or ""
        if p and Path(p).exists():
            return p
    return None


# ── Cover slide ───────────────────────────────────────────────────────────────


def _cover(prs: Presentation, brand: str, hero_path: str | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = SLIDE_W, SLIDE_H

    _fill_white(slide)

    panel_w = Inches(6.15)
    gap = Inches(0.22)

    # Left — hero image or soft panel
    if hero_path:
        _add_picture_cover(slide, hero_path, Inches(0), Inches(0), panel_w, H)
        # Light gradient overlay strip on the right edge of the photo for separation
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, panel_w - Inches(0.08), 0, Inches(0.12), H)
        ov.fill.solid()
        ov.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        ov.fill.transparency = 0.65
        ov.line.fill.background()
    else:
        _rect(slide, 0, 0, panel_w, H, COVER_PANEL)
        _rect_outline(slide, gap, gap, panel_w - 2 * gap, H - 2 * gap, RULE, 0.75)

    right_x = panel_w + Inches(0.55)

    _textbox(
        slide,
        text=brand.upper(),
        x=right_x,
        y=Inches(2.65),
        w=Inches(6.4),
        h=Inches(1.0),
        size=34,
        bold=False,
        color=BLACK,
        align=PP_ALIGN.LEFT,
    )

    _rect(slide, right_x, Inches(3.72), Inches(5.85), Inches(0.02), ACCENT)

    _textbox(
        slide,
        text="AW 2026  ·  STYLE SELECTION",
        x=right_x,
        y=Inches(3.92),
        w=Inches(6.4),
        h=Inches(0.45),
        size=11,
        bold=False,
        color=MUTED,
        align=PP_ALIGN.LEFT,
        letter_spacing=120,
    )

    _textbox(
        slide,
        text="STYLE SPECIFICATIONS & TECHNICAL DATA",
        x=right_x,
        y=Inches(4.52),
        w=Inches(6.4),
        h=Inches(0.4),
        size=9,
        bold=False,
        color=FAINT,
        align=PP_ALIGN.LEFT,
        letter_spacing=200,
    )


# ── Product slide — matches StyleGroupCluster (header + spec + 2-col grid) ────


def _product_slide(prs: Presentation, group: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = SLIDE_W, SLIDE_H

    _fill_white(slide)

    by_type: dict[str, list[str]] = {
        "front": [],
        "back": [],
        "detail": [],
        "spec_label": [],
    }
    for img in group.get("images", []) or []:
        t = img.get("image_type", "unknown")
        src = img.get("processed_path") or img.get("original_path", "")
        if src and Path(src).exists():
            bucket = by_type.get(t)
            if bucket is not None:
                bucket.append(src)

    # Same ordering as the UI card: front → back → all details → spec label
    ordered = (
        by_type["front"][:1]
        + by_type["back"][:1]
        + by_type["detail"]
        + by_type["spec_label"][:1]
    )

    pad = Inches(0.38)
    header_h = Inches(0.62)
    gap_y = Inches(0.14)
    gap_x = Inches(0.12)

    style_id = (group.get("style_id") or "").strip() or "—"
    gtype = (group.get("garment_type") or "").strip()
    gdata = group.get("garment_data") or {}

    ref = (gdata.get("reference_number") or "").strip()
    comp = (gdata.get("fabric_composition") or "").strip()
    gsm = (gdata.get("gsm") or "").strip()
    date = (gdata.get("date") or "").strip()

    has_spec = bool(ref or comp or gsm or date)

    spec_h = Inches(1.05) if has_spec else Inches(0)
    if has_spec:
        spec_h = _spec_panel_height(ref, comp, gsm, date)

    y_cursor = pad

    # ── Header row (style id + garment type) ─────────────────────────────────
    _slide_header(slide, pad, y_cursor, W - 2 * pad, header_h, style_id, gtype)
    y_cursor += header_h + gap_y

    # ── Spec panel ───────────────────────────────────────────────────────────
    if has_spec:
        _spec_panel(slide, pad, y_cursor, W - 2 * pad, spec_h, ref, comp, gsm, date)
        y_cursor += spec_h + gap_y

    # ── Image grid ───────────────────────────────────────────────────────────
    grid_top = y_cursor
    grid_bottom = H - pad - Inches(0.42)
    grid_h = grid_bottom - grid_top
    grid_w = W - 2 * pad

    n = len(ordered)

    if n == 0:
        _place_placeholder(slide, pad, grid_top, grid_w, grid_h, "NO IMAGES")
    else:
        slots = _grid_slots(n, pad, grid_top, grid_w, grid_h, gap_x, gap_y)
        for img_path, (sx, sy, sw, sh) in zip(ordered, slots):
            # Thin frame per thumbnail (like UI card border)
            frame_pad = Inches(0.02)
            inner_x, inner_y = sx + frame_pad, sy + frame_pad
            inner_w, inner_h = sw - 2 * frame_pad, sh - 2 * frame_pad
            border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, sw, sh)
            border.fill.background()
            border.line.color.rgb = RULE
            border.line.width = Pt(0.75)
            try:
                border.adjustments[0] = 0.06
            except Exception:
                pass
            _add_picture_cover(slide, img_path, inner_x, inner_y, inner_w, inner_h)

    # ── Slide index (frontend: SLIDE N chip) ────────────────────────────────
    _textbox(
        slide,
        text=f"SLIDE {slide_num:02d}",
        x=W - Inches(1.35),
        y=H - Inches(0.34),
        w=Inches(1.2),
        h=Inches(0.26),
        size=8,
        bold=False,
        color=FAINT,
        align=PP_ALIGN.RIGHT,
        font_name="Calibri",
    )


def _slide_header(slide, x: Emu, y: Emu, w: Emu, h: Emu, style_id: str, garment_type: str):
    """Title row: bold style id (left), garment type (right)."""
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.margin_bottom = tf.margin_top = tf.margin_left = tf.margin_right = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    r1 = p.add_run()
    r1.text = style_id
    r1.font.bold = True
    r1.font.size = Pt(17)
    r1.font.color.rgb = BLACK
    r1.font.name = "Calibri"

    if garment_type:
        r2 = p.add_run()
        r2.text = f"     ·     {garment_type.upper()}"
        r2.font.bold = False
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED
        r2.font.name = "Calibri"


def _spec_panel_height(ref: str, comp: str, gsm: str, date: str) -> Emu:
    rows = sum(1 for v in (ref, comp, gsm, date) if v)
    base = Inches(0.38)
    per_row = Inches(0.21)
    return base + per_row * max(rows, 1)


def _spec_panel(
    slide,
    x: Emu,
    y: Emu,
    w: Emu,
    h: Emu,
    ref: str,
    comp: str,
    gsm: str,
    date: str,
):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL_BG
    panel.line.color.rgb = PANEL_BORDER
    panel.line.width = Pt(0.5)
    try:
        panel.adjustments[0] = 0.04
    except Exception:
        pass

    inner = Inches(0.14)
    tx = slide.shapes.add_textbox(x + inner, y + inner * 0.85, w - 2 * inner, h - 1.6 * inner)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.auto_size = MSO_AUTO_SIZE.NONE

    rows: list[tuple[str, str]] = []
    if ref:
        rows.append(("REF", ref))
    if comp:
        rows.append(("COMP", comp))
    if gsm:
        rows.append(("GSM", gsm))
    if date:
        rows.append(("DATE", date))

    for i, (lab, val) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else 3)
        p.space_after = Pt(0)
        p.line_spacing = 1.05

        rl = p.add_run()
        rl.text = f"{lab}  "
        rl.font.bold = True
        rl.font.size = Pt(9)
        rl.font.color.rgb = FAINT
        rl.font.name = "Calibri"

        rv = p.add_run()
        rv.text = val[:120]
        rv.font.bold = False
        rv.font.size = Pt(10)
        rv.font.color.rgb = SUBTITLE
        rv.font.name = "Calibri"


def _grid_slots(
    n: int,
    gx: Emu,
    gy: Emu,
    gw: Emu,
    gh: Emu,
    gap_x: Emu,
    gap_y: Emu,
) -> list[tuple[Emu, Emu, Emu, Emu]]:
    """Two-column slots; single image is centred with portrait-friendly 3:4 box."""
    if n <= 0:
        return []
    if n == 1:
        gw_i, gh_i = _emu(gw), _emu(gh)
        # Target portrait 3:4 inside grid
        slot_w = min(gw_i, int(gh_i * 3 / 4))
        slot_h = int(slot_w * 4 / 3)
        if slot_h > gh_i:
            slot_h = gh_i
            slot_w = int(slot_h * 3 / 4)
        ox = _emu(gx) + (gw_i - slot_w) // 2
        oy = _emu(gy) + (gh_i - slot_h) // 2
        return [(Emu(ox), Emu(oy), Emu(slot_w), Emu(slot_h))]

    cols = 2
    rows = (n + cols - 1) // cols
    cw = (_emu(gw) - _emu(gap_x)) // 2
    ch = (_emu(gh) - _emu(gap_y) * max(rows - 1, 0)) // max(rows, 1)

    out: list[tuple[Emu, Emu, Emu, Emu]] = []
    for i in range(n):
        r, c = divmod(i, cols)
        sx = _emu(gx) + c * (cw + _emu(gap_x))
        sy = _emu(gy) + r * (ch + _emu(gap_y))
        out.append((Emu(sx), Emu(sy), Emu(cw), Emu(ch)))
    return out


# ── Closing slide ─────────────────────────────────────────────────────────────


def _closing(prs: Presentation, brand: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = SLIDE_W, SLIDE_H

    _fill_white(slide)

    rule_y = Inches(3.17)
    _rect(slide, Inches(1.05), rule_y, Inches(4.17), Inches(0.018), RULE)
    _rect(slide, Inches(7.99), rule_y, Inches(4.17), Inches(0.018), RULE)

    _textbox(
        slide,
        text="www.garmentsupply.com",
        x=Inches(5.35),
        y=Inches(2.85),
        w=Inches(2.65),
        h=Inches(0.42),
        size=12,
        bold=False,
        color=BLACK,
        align=PP_ALIGN.CENTER,
    )

    _textbox(
        slide,
        text="Bengaluru  ·  Chennai  ·  Gurugram  ·  Tirupur",
        x=Inches(4.05),
        y=Inches(4.51),
        w=Inches(5.25),
        h=Inches(0.38),
        size=11,
        bold=False,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    _textbox(
        slide,
        text=f"{brand.upper()}  ·  2026",
        x=Inches(4.05),
        y=H - Inches(0.62),
        w=Inches(5.25),
        h=Inches(0.38),
        size=9,
        bold=False,
        color=FAINT,
        align=PP_ALIGN.CENTER,
        letter_spacing=100,
    )


def _place_placeholder(slide, x, y, w, h, label: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_BG
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(10)
    run.font.color.rgb = FAINT


def _fill_white(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PAGE_BG


def _rect(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _rect_outline(slide, x, y, w, h, color: RGBColor, width_pt: float = 0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)
    return shape


def _textbox(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    size=12,
    bold=False,
    color: RGBColor = BLACK,
    align=PP_ALIGN.LEFT,
    letter_spacing: int = 0,
    font_name: str = "Calibri",
):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font_name

    if letter_spacing:
        r_pr = run._r.get_or_add_rPr()
        r_pr.set("spc", str(letter_spacing))

    return txb
